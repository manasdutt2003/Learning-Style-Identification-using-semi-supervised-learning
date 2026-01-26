from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def check_file_exists(filename):
    if not os.path.exists(filename):
        print(f"Warning: {filename} not found.")
        return False
    return True

def add_slide(prs, layout_index, title_text, content_text=None):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    if content_text:
        # Check if placeholders exist
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            content.text = content_text
    
    return slide

def apply_formatting(slide):
    # Basic formatting to make it look "Ivy League" (Clean, Serif fonts usually)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = 'Garamond' # Or Arial for safety
            if paragraph.font.size is None:
                paragraph.font.size = Pt(18)

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Driven Learning Style Prediction for Adaptive E-Learning Systems"
    subtitle.text = "Team Members:\nManas Dutt (22BCE1350)\nUrvi Samirkumar Shah (22BCE1174)\nSoham Yogesh Amberkar (22BCE1770)\n\nGuide: Dr. Vijayalakshmi A"

    # --- Slide 2: Introduction ---
    content = ("• E-Learning (MOOCs, LMS) is growing but often lacks personalization.\n"
               "• One-size-fits-all approach reduces student engagement.\n"
               "• Traditional learning style identification relies on static surveys.\n"
               "• Need for real-time, behavior-based detection using AI.")
    slide = add_slide(prs, 1, "Introduction", content) # 1 = Title and Content

    # --- Slide 3: Problem Statement ---
    content = ("• Existing LMS Lack Scalability: Manual surveys are impractical for large classes.\n"
               "• Static & Biased: Self-reported questionnaires may not reflect actual behavior.\n"
               "• Underuse of Advanced AI: Modern techniques like Graph Learning are unexplored.\n"
               "• Lack of Explainability: Educators need to understand 'WHY' a style is predicted.")
    add_slide(prs, 1, "Problem Statement", content)

    # --- Slide 4: Objectives ---
    content = ("• Develop an automated, AI-driven system to predict learning styles from LMS logs.\n"
               "• Implement Semi-Supervised Learning to handle limited labeled data.\n"
               "• Map behavioral features to the Felder-Silverman Learning Style Model (FSLSM).\n"
               "• Create an Analytics Dashboard for faculty insights.")
    add_slide(prs, 1, "Objectives", content)

    # --- Slide 5: Methodology (With Architecture) ---
    slide = add_slide(prs, 5, "Proposed Methodology") # 5 = Title Only
    
    # Add text box on the left
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "1. Data Collection: Extract logs (video views, quizzes).\n"
    p = tf.add_paragraph()
    p.text = "2. Feature Engineering: Align actions with FSLSM dimensions."
    p = tf.add_paragraph()
    p.text = "3. Modeling: Hybrid Neural Networks & Ensemble Boosting."
    p = tf.add_paragraph()
    p.text = "4. Semi-Supervised: Leverage unlabeled data."
    
    # Add Architecture Image on the right
    img_path = 'architecture_diagram.png'
    if check_file_exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5), Inches(1.5), width=Inches(4.5))

    # --- Slide 6: System Architecture ---
    # Dedicated slide for the diagram to be large
    slide = add_slide(prs, 5, "System Architecture")
    if check_file_exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

    # --- Slide 7: Tech Stack ---
    content = ("• Languages: Python (3.x)\n"
               "• Machine Learning: Scikit-learn, LightGBM, CatBoost, PyTorch/TensorFlow\n"
               "• Data Processing: Pandas, NumPy, SMOTE (Imbalanced-learn)\n"
               "• Visualization: Matplotlib, Seaborn, SHAP (Explainability)\n"
               "• Tools: Jupyter Notebooks, Git")
    add_slide(prs, 1, "Tools and Technologies", content)

    # --- Slide 8: Expected Outcomes ---
    content = ("• >85% Accuracy in Learning Style Prediction.\n"
               "• Real-time, scalable detection without surveys.\n"
               "• Actionable insights for instructors via dashboard.\n"
               "• Framework for future Adaptive Content Delivery.")
    add_slide(prs, 1, "Expected Outcomes", content)
    
    # --- Slide 9: References ---
    content = ("1. IEEE Transactions on Learning Technologies (2024). 'Learning Style Identification...'\n"
               "2. Felder, R. M. & Silverman, L. K. (1988). 'Learning and Teaching Styles in Engineering Education.'\n"
               "3. GitHub Datasets: 'Learning-style-prediction-identification-using-Machine-Learning'")
    slide = add_slide(prs, 1, "References", content)
    
    # Resize text for references
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(14)

    # --- Slide 10: Conclusion ---
    slide = add_slide(prs, 1, "Thank You", "")
    txBox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(2))
    tf = txBox.text_frame
    tf.text = "Questions?"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(40)

    save_path = "Research_Project_Presentation.pptx"
    prs.save(save_path)
    print(f"Presentation saved successfully to {save_path}")

if __name__ == "__main__":
    create_presentation()
